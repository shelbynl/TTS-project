"""
Extract Presenter Notes
TODO: #2 add menu to main function
TODO: #3 add location to save text to
TODO: #4 add method for multiple files
TODO: #5 make FILE an input with type checking

REFERENCE:
<https://stackoverflow.com/questions/63659972/extract-presenter-notes-from-pptx-file-powerpoint>
"""

import os
from pathlib import Path
import collections.abc
from pptx import Presentation


def main():
    """runs the main function of the program"""
    ret_path = file_names()
    ppt = Presentation(ret_path)
    lst = enum_notes(ppt)
    with open(file="note.csv", mode="w", encoding="utf-8") as file_out:
        path = os.getcwd(lst)
        for tup in path:
            tup.write(lst[1::])
        return file_out



def file_names():
    """
    Function that takes in the file names from the user
       and checks if they have the correct endings
    """
    print("Please enter a .pptx file for audio translation\n" +
          "Enter 'q' to quit.")
    file1 = input("\nFile1: ")
    path = Path(file1)
    if path == str('q'):
        exit()
    return path

    # if os.path.exists(p1):
    #     f = p1
    #     return f
    # try:
    #     answer = int(first_number) / int(second_number)
    # except TypeError():
    #     print("You specified an incorrect file ending")
    # else:
    #     print(answer)


def enum_notes(pptx_scrape):
    """Enumerates the notes and returns the text"""
    notes = []
    for page, slide in enumerate(pptx_scrape.slides):
        text_note = slide.notes_slide.notes_text_frame.text
        notes.append((page, text_note))
    return notes


def enum_slides(pptx_scrape):
    """Enumerates the slides for text"""
    notes = []
    for page, slide in enumerate(pptx_scrape.slides):
        temp = []
        for shape in slide.shapes:
            # this will extract all text in text boxes on the slide.
            if shape.has_text_frame and shape.text.strip():
                temp.append(shape.text)
        notes.append((page, temp))
    return notes


if __name__ == "__main__":
    main()
