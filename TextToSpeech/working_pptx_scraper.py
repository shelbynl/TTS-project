"""
Extract Presenter Notes
TODO: #2 add menu to main function
TODO: #3 add location to save text to
TODO: #4 add method for multiple files
TODO: #5 make FILE an input with type checking

REFERENCE:
<https://stackoverflow.com/questions/63659972/extract-presenter-notes-from-pptx-file-powerpoint>
"""

import collections.abc
from pptx import Presentation
import pandas as pd
import numpy as np
import pyinputplus as pyip


# file_in = input("Input powerpoint file location for audio:\n")
# dir_out = input("Input directory for output files:\n")

FILE_IN = "D:\\Repositories\\TTS\\CareerNavigation\\New-Slides\\07-CS9101.pptx"
FILE_OUT = "D:\\Repositories\\TTS\\CareerNavigation\\output_text\\07.csv"
PPT = Presentation(FILE_IN)


def main():
    """runs the main function of the program"""
    menu()
    scraped_text = enum_notes()
    file_o(scraped_text)


def menu():
    """
    menu creates the function that will take
    the input from the user and validate it
    """
    result = pyip.inputMenu(['Filename in',''], lettered=True)


    return result

def enum_notes():
    """
    Enumerates the notes and returns the text
    """
    notes = []
    for page, slide in enumerate(PPT.slides):
        txt_note = slide.notes_slide.notes_text_frame.text
        notes.append((page, txt_note))
    return notes


def enum_slides():
    """
    Enumerates the slides for text
    """
    notes = []
    for page, slide in enumerate(PPT.slides):
        temp = []
        for shape in slide.shapes:
            # this will extract all text in text boxes on the slide.
            if shape.has_text_frame and shape.text.strip():
                temp.append(shape.text)
        notes.append((page, temp))
    return notes


def file_o(new_lst):
    """
    The file it will be written to
    """
    with open(file=FILE_OUT, mode='w', encoding='utf-8') as tmp_file:
        data = pd.DataFrame(new_lst)
        data.to_csv(tmp_file, columns=data[:-1])


if __name__=='__main__':
    main()
