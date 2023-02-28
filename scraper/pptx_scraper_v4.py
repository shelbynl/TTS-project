"""PPTX_Scraper_v4.py

Changes from version 3

Here are the main changes I made:

I created a class called PPTX_Scraper that
    encapsulates all the functions in the file.

I moved the __init__ function and the main function into the class
    and modified them to use the class's instance variables.

I modified the file_o function to take the notes argument
    and write each note to a separate file with a filename
    of the form page_X.txt, where X is the page number.

I modified the walk function to also use self and made it an

Walk() Function referenced from
https://learning.oreilly.com/library/view/the-recursive-book/9781098141356/c10.xhtml#h1-502024c10-0001
"""

import os
import pandas as pd
from pptx import Presentation


class PPTX_Scraper:
    def __init__(self):
        self.file_in = input("Input powerpoint file location for audio:\n")
        self.dir_out = input("Input directory for output files:\n")

    def main(self):
        """runs the main function of the program"""
        ppt = Presentation(self.file_in)
        notes = self.enum_notes(ppt)
        self.file_o(notes)

    def enum_notes(self, ppt):
        """Enumerates the notes and returns the text"""
        notes = []
        for page, slide in enumerate(ppt.slides):
            txt_note = slide.notes_slide.notes_text_frame.text
            notes.append((page, txt_note))
        return notes

    def enum_slides(self, ppt):
        """Enumerates the slides for text"""
        notes = []
        for page, slide in enumerate(ppt.slides):
            temp = []
            for shape in slide.shapes:
                # this will extract all text in text boxes on the slide.
                if shape.has_text_frame and shape.text.strip():
                    temp.append(shape.text)
            notes.append((page, temp))
        return notes

    def file_o(self, notes):
        """
        File that will be written out to.
        """
        for page, note in notes:
            file_name = f"page_{page}.txt"
            with open(os.path.join(self.dir_out, file_name), "w", encoding="utf-8") as f:
                f.write(note)


    def walk(self, folder, matchFunc):
        """Calls the match function with every file in the folder and its
        subfolders. Returns a list of files that the match function
        returned True for."""
        matchedFiles = [] # This list holds all the matches.
        folder = os.path.abspath(folder) # Use the folder's absolute path.

        # Loop over every file and subfolder in the folder:
        for name in os.listdir(folder):
            filepath = os.path.join(folder, name)
            if os.path.isfile(filepath):
                # Call the match function for each file:
                if matchFunc(filepath):
                    matchedFiles.append(filepath)
            elif os.path.isdir(filepath):
                # Recursively call walk for each subfolder, extending
                # the matchedFiles with their matches:
                matchedFiles.extend(self.walk(filepath, matchFunc))
        return matchedFiles


if __name__=='__main__':
    scraper = PPTX_Scraper()
    option = input("Would you like to scrape a single file or multiple files? (s/m): ")
    if option == "s":
        scraper.main()
    elif option == "m":
        folder = input("Input the directory to search for PowerPoint files: ")
        ppt_files = scraper.walk(folder, lambda x: x.endswith(".pptx"))
        for ppt_file in ppt_files:
            scraper.file_in = ppt_file
            scraper.main()
    else:
        print("Invalid option. Please choose 's' or 'm'.")
