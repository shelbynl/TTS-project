"""
Extract Presenter Notes
TODO: #2 add menu to main function
TODO: #3 add location to save text to
TODO: #4 add method for multiple files
TODO: #5 make FILE an input with type checking

REFERENCE:
<https://stackoverflow.com/questions/63659972/extract-presenter-notes-from-pptx-file-powerpoint>
"""

import collections.abc
import os
import pandas as pd
from pptx import Presentation
from pathlib import Path



def __init__(self):
    self.file_in = input("Input powerpoint file location for audio:\n")
    self.dir_out = input("Input directory for output files:\n")

def main(self):
    """runs the main function of the program"""
    self.ppt = Presentation(self.file_in)
    enum_notes(self)
    file_o(self.dir_out)


def enum_notes(self):
    """Enumerates the notes and returns the text"""
    notes = []
    for page, slide in enumerate(self.ppt.slides):
        txt_note = slide.notes_slide.notes_text_frame.text
        notes.append((page, txt_note))
    return notes


def enum_slides(self):
    """Enumerates the slides for text"""
    notes = []
    for page, slide in enumerate(self.ppt.slides):
        temp = []
        for shape in slide.shapes:
            # this will extract all text in text boxes on the slide.
            if shape.has_text_frame and shape.text.strip():
                temp.append(shape.text)
        notes.append((page, temp))
    return notes


def file_o(outfile):
    """
    and the file it will be written to
    """
    new_lst = []
    with open(file=outfile, mode='w', encoding='utf-8') as file_out:
        df = pd.DataFrame(new_lst)
        temp = df.to_csv(file_out, index=False, header=False)
#        for i[0:2] in temp:
#            file_name =f'{i[1]}.txt'
#            with open(file_name, 'w', encoding='utf-8') as f:
#                f.write(i['Seq'])


def walk(folder, matchFunc):
    """Calls the match function with every file in the folder and its
    subfolders. Returns a list of files that the match function
    returned True for."""
    matchedFiles = [] # This list holds all the matches.
    folder = os.path.abspath(folder) # Use the folder's absolute path.

    # Loop over every file and subfolder in the folder:
    for name in os.listdir(folder):
        filepath = os.path.join(folder, name)
        if os.path.isfile(filepath):
            # Call the match function for each file:
            if matchFunc(filepath):
                matchedFiles.append(filepath)
        elif os.path.isdir(filepath):
            # Recursively call walk for each subfolder, extending
            # the matchedFiles with their matches:
            matchedFiles.extend(walk(filepath, matchFunc))
    return matchedFiles


if __name__=='__main__':
    main()
